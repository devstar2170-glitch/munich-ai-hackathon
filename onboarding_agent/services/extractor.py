"""Gemini-based CV data extraction."""

import json
import logging

import google.generativeai as genai

from config import GEMINI_API_KEY, GEMINI_MODEL
from models import CandidateProfile, SupplementExtraction

logger = logging.getLogger(__name__)

SUPPLEMENT_SYSTEM_INSTRUCTION = """You are reviewing a supplementary document (CV update, project \
deck, certificate, etc.) uploaded by an employee to fill gaps in their staffing profile.

Rules:
- Respond with ONLY a single valid JSON object. No markdown, no commentary, no code fences.
- Only return fields from the "Gap fields" list below — ignore anything else, even if mentioned.
- Only include a field if the document gives you genuinely new, specific information about it.
- "value" must always be a string. For list fields (e.g. skills, certifications,
  pastIndustryExperience, futureIndustryWish), comma-separate the items in a single string.
- "confidence" is 0-100: how confident you are this value is correct and ready to apply directly
  to the profile without human review. Use lower confidence for vague, inferred, or ambiguous
  information.
- "reasoning" is a one-sentence justification referencing what in the document supports the value.
"""

SYSTEM_INSTRUCTION = """You are a CV/resume parser for a recruitment CRM.
Given the contents of a candidate's CV (PDF), extract the requested fields.

Rules:
- Respond with ONLY a single valid JSON object. No markdown, no commentary, no code fences.
- If a field cannot be determined from the CV, use null (or an empty list for list fields).
- "level" must be exactly one of: "Junior", "Mid", "Senior", "Lead".
- "yearsOfExperience" must be a number (total years of professional experience).
- "skills" should include technical and relevant soft skills mentioned in the CV.
- "pastIndustryExperience" should list the industries the candidate has worked in.
"""

_model = None


def _get_model() -> genai.GenerativeModel:
    global _model
    if _model is None:
        if not GEMINI_API_KEY:
            raise RuntimeError("GEMINI_API_KEY environment variable is not set.")
        genai.configure(api_key=GEMINI_API_KEY)
        _model = genai.GenerativeModel(
            model_name=GEMINI_MODEL,
            system_instruction=SYSTEM_INSTRUCTION,
            generation_config={
                "response_mime_type": "application/json",
                "response_schema": CandidateProfile,
            },
        )
    return _model


def extract_candidate_profile(pdf_bytes: bytes, file_name: str) -> CandidateProfile:
    """Send a CV PDF to Gemini and return a validated CandidateProfile."""
    model = _get_model()

    logger.info("Sending '%s' to Gemini (%s) for extraction...", file_name, GEMINI_MODEL)
    response = model.generate_content(
        [
            {"mime_type": "application/pdf", "data": pdf_bytes},
            "Extract the candidate's profile information from this CV.",
        ]
    )

    raw_text = response.text
    try:
        data = json.loads(raw_text)
    except json.JSONDecodeError as exc:
        raise ValueError(f"Gemini did not return valid JSON for '{file_name}': {raw_text!r}") from exc

    return CandidateProfile.model_validate(data)


def extract_supplement_fields(
    file_bytes: bytes,
    mime_type: str,
    file_name: str,
    profile: dict,
    gap_fields: list[str],
) -> SupplementExtraction:
    """Send a supplementary document plus the current profile to Gemini and return
    candidate values + confidence scores for the given gap fields only."""
    if not GEMINI_API_KEY:
        raise RuntimeError("GEMINI_API_KEY environment variable is not set.")

    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel(
        model_name=GEMINI_MODEL,
        system_instruction=SUPPLEMENT_SYSTEM_INSTRUCTION,
        generation_config={
            "response_mime_type": "application/json",
            "response_schema": SupplementExtraction,
        },
    )

    prompt = f"""
Current profile:
{json.dumps(profile, indent=2, ensure_ascii=False)}

Gap fields (only consider these):
{json.dumps(gap_fields, indent=2, ensure_ascii=False)}
"""

    logger.info("Sending '%s' to Gemini (%s) for supplement extraction...", file_name, GEMINI_MODEL)
    response = model.generate_content([
        {"mime_type": mime_type, "data": file_bytes},
        prompt,
    ])

    raw_text = response.text
    try:
        data = json.loads(raw_text)
    except json.JSONDecodeError as exc:
        raise ValueError(f"Gemini did not return valid JSON for '{file_name}': {raw_text!r}") from exc

    return SupplementExtraction.model_validate(data)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract all text from a PPTX file's slides."""
    from io import BytesIO

    from pptx import Presentation

    presentation = Presentation(BytesIO(file_bytes))
    chunks = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        slide_chunks.append(text)
        if slide_chunks:
            chunks.append(f"Slide {slide_index}:\n" + "\n".join(slide_chunks))
    return "\n\n".join(chunks)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract all text from a DOCX file's paragraphs."""
    from io import BytesIO

    from docx import Document

    document = Document(BytesIO(file_bytes))
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())
